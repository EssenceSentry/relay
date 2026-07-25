from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Protocol, cast

from knowledge_core.models import TextSection


class _TextShape(Protocol):
    @property
    def text(self) -> str: ...


class UnsupportedDocumentError(ValueError):
    pass


def parse_document(data: bytes, filename: str) -> list[TextSection]:
    suffix = Path(filename).suffix.casefold()
    if suffix == ".pdf":
        return _parse_pdf(data)
    if suffix == ".docx":
        return _parse_docx(data)
    if suffix == ".doc":
        return _parse_doc(data)
    if suffix == ".pptx":
        return _parse_pptx(data)
    if suffix in {".txt", ".md", ".csv", ".log"}:
        return _parse_plain_text(data)
    if suffix == ".json":
        return _parse_json(data)
    raise UnsupportedDocumentError(
        f"Unsupported document type {suffix or '<none>'}. "
        "Supported: PDF, DOC, DOCX, PPTX, TXT, MD, CSV, JSON."
    )


def _decode(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _parse_plain_text(data: bytes) -> list[TextSection]:
    text = _decode(data).strip()
    return [TextSection(text=text, locator="document")] if text else []


def _parse_json(data: bytes) -> list[TextSection]:
    raw = _decode(data)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return _parse_plain_text(data)
    text = json.dumps(parsed, ensure_ascii=False, indent=2)
    return [TextSection(text=text, locator="document")]


def _parse_pdf(data: bytes) -> list[TextSection]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    sections: list[TextSection] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            sections.append(
                TextSection(
                    text=text,
                    locator=f"page {index}",
                    title=f"Page {index}",
                    page_number=index,
                )
            )
    return sections


def _parse_docx(data: bytes) -> list[TextSection]:
    from docx import Document

    document = Document(io.BytesIO(data))
    sections: list[TextSection] = []
    current_title: str | None = None
    current_lines: list[str] = []

    def flush() -> None:
        nonlocal current_lines
        text = "\n\n".join(line for line in current_lines if line).strip()
        if text:
            sections.append(
                TextSection(
                    text=text,
                    title=current_title,
                    locator=current_title or "document",
                )
            )
        current_lines = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue
        style_name = (
            (paragraph.style.name or "").casefold() if paragraph.style else ""
        )
        if style_name.startswith("heading"):
            flush()
            current_title = text
        else:
            current_lines.append(text)

    for table_index, table in enumerate(document.tables, start=1):
        rows: list[str] = []
        for row in table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                rows.append(" | ".join(values))
        if rows:
            current_lines.append(f"Table {table_index}\n" + "\n".join(rows))
    flush()
    return sections


def _parse_doc(data: bytes) -> list[TextSection]:
    import unword

    document = unword.parse_doc(data)
    sections: list[TextSection] = []
    body = str(document.body_text or "").strip()
    if body:
        sections.append(TextSection(text=body, locator="document"))

    for index, raw_text in enumerate(document.textboxes, start=1):
        text = str(raw_text or "").strip()
        if text:
            sections.append(
                TextSection(
                    text=text,
                    title=f"Text box {index}",
                    locator=f"text box {index}",
                )
            )
    return sections


def _parse_pptx(data: bytes) -> list[TextSection]:
    from pptx import Presentation
    from pptx.shapes.graphfrm import GraphicFrame

    presentation = Presentation(io.BytesIO(data))
    sections: list[TextSection] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines: list[str] = []
        title: str | None = None
        if slide.shapes.title is not None:
            title = slide.shapes.title.text.strip() or None

        for shape in slide.shapes:
            if shape.has_text_frame:
                text = cast(_TextShape, shape).text.strip()
                if text:
                    lines.append(text)
            if isinstance(shape, GraphicFrame) and shape.has_table:
                for row in shape.table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        lines.append(" | ".join(values))

        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame is not None:
            notes = notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append(f"Speaker notes:\n{notes}")

        text = "\n\n".join(dict.fromkeys(lines)).strip()
        if text:
            sections.append(
                TextSection(
                    text=text,
                    locator=f"slide {slide_index}",
                    title=title or f"Slide {slide_index}",
                    page_number=slide_index,
                )
            )
    return sections
